#!/usr/bin/env python3
"""
Generate NHS Safeguarding Presentation PowerPoint
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

# NHS Blue color
NHS_BLUE = RGBColor(0, 94, 184)
NHS_DARK_BLUE = RGBColor(0, 48, 135)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
RED = RGBColor(218, 41, 28)
AMBER = RGBColor(255, 191, 0)
GREEN = RGBColor(0, 135, 68)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Blue background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = NHS_BLUE
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, highlight_color=None):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = NHS_BLUE
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(11.9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = bullet
        p.font.size = Pt(24)
        p.font.color.rgb = BLACK
        p.level = 0
        p.space_before = Pt(12)
        p.space_after = Pt(6)
    
    return slide

def add_two_column_slide(prs, title, left_content, right_content):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = NHS_BLUE
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = BLACK
        p.space_before = Pt(8)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = BLACK
        p.space_before = Pt(8)
    
    return slide

# === CREATE SLIDES ===

# Slide 1: Title
add_title_slide(prs, 
    "Radio Check", 
    "AI Safeguarding Framework for Veteran Mental Health\nSeeking NHS Guidance & Partnership")

# Slide 2: The Challenge
add_content_slide(prs, "The Challenge: UK Veteran Mental Health Crisis", [
    "• 13% of UK veterans experience PTSD (vs 4% general population)",
    "• Veterans are 2-3x more likely to experience depression",
    "• High rates of isolation, especially in rural areas",
    "• Stigma prevents many from seeking traditional help",
    "• Average of 2 veteran suicides per week in the UK",
    "",
    "\"Many veterans won't walk into a clinic, but they'll open an app at 3am\""
])

# Slide 3: Radio Check Solution
add_content_slide(prs, "Radio Check: The Solution", [
    "A mobile-first mental health support app for UK veterans",
    "",
    "• AI 'Battle Buddies' - companions trained in military culture",
    "• 24/7 availability - support when veterans need it most",
    "• No appointments, no waiting lists",
    "• Familiar, non-clinical environment",
    "• Direct connection to human support when needed",
    "",
    "NOT a replacement for clinical care - a bridge to it"
])

# Slide 4: The Safeguarding Challenge
add_content_slide(prs, "The Safeguarding Challenge", [
    "When AI Meets Crisis:",
    "",
    "• Users may disclose suicidal thoughts to AI before humans",
    "• Subtle language cues can indicate severe distress",
    "• Typos and informal language in crisis moments",
    "• Need for immediate, appropriate response",
    "• Balance between support and escalation",
    "",
    "Our Question: How do we build AI that saves lives?"
])

# Slide 5: 4-Level Risk Model
add_content_slide(prs, "Our 4-Level Risk Detection Model", [
    "LEVEL 0: Normal conversation → Continue supportive chat",
    "",
    "LEVEL 1: Low distress (sadness, stress) → Acknowledge, invite to talk",
    "",
    "LEVEL 2: Hopelessness language → SAFEGUARDING BEGINS",
    "",
    "LEVEL 3: Self-harm thoughts → Safety check, provide resources",
    "",
    "LEVEL 4: Imminent risk → CRISIS MODE, immediate resources",
    "",
    "Key Principle: Don't wait for explicit suicide statements"
])

# Slide 6: What We Detect
add_two_column_slide(prs, "What We Detect", 
    [
        "308 RED Indicators",
        "(Immediate Escalation)",
        "",
        "• \"I won't be here tomorrow\"",
        "• \"I can't go on anymore\"",
        "• \"Everyone would be better without me\"",
        "• \"I've made a plan\"",
        "",
        "Including temporal variations:",
        "tonight, tomorrow, in the morning, soon"
    ],
    [
        "255 AMBER Indicators",
        "(Early Warning)",
        "",
        "• \"There's no point\"",
        "• \"I'm trapped\"",
        "• \"Nobody understands\"",
        "• \"Drinking to cope\"",
        "",
        "PTSD triggers:",
        "flashbacks, nightmares, hypervigilance"
    ])

# Slide 7: Real Example
add_content_slide(prs, "Handling Real-World Input", [
    "Users in crisis type quickly with mistakes. We detect anyway:",
    "",
    "User types: \"cant see me beeing here in the morning\"",
    "We detect: \"can't see me being here in the morning\" → RED FLAG",
    "",
    "Example Conversation That Now Triggers RED Alert:",
    "",
    "\"hey feeking down not sureif i can go on\"",
    "\"just no point in going on\"",
    "\"too much I cant see me beeing here in the morning\"",
    "",
    "All three messages = IMMEDIATE RED ALERT + Crisis Resources"
])

# Slide 8: AI Response Protocol
add_content_slide(prs, "AI Response Protocol", [
    "Level 2-3 (Hopelessness / Self-harm thoughts):",
    "• \"I'm really glad you reached out.\"",
    "• \"That sounds incredibly heavy to carry.\"",
    "• \"Are you safe where you are right now?\"",
    "",
    "Level 4 (Imminent Risk) - Crisis Resources:",
    "• NHS Mental Health: 111 (Option 2) - Free, 24/7",
    "• Samaritans: 116 123 - Free, 24/7",
    "• Text SHOUT: 85258 - Free, 24/7",
    "• Emergency: 999",
    "",
    "\"I'm here with you. You don't have to face this alone.\""
])

# Slide 9: What We Need From NHS
add_content_slide(prs, "What We Need From NHS", [
    "1. Clinical Review",
    "   - Are our trigger phrases clinically appropriate?",
    "   - Are we missing critical indicators?",
    "",
    "2. Response Guidance",
    "   - Are AI responses clinically sound?",
    "   - What language works best for de-escalation?",
    "",
    "3. Integration Opportunities",
    "   - Connection to Op Courage pathways",
    "   - NHS 111 Option 2 referral protocols",
    "",
    "4. Ongoing Partnership",
    "   - Regular clinical review and refinement"
])

# Slide 10: Our Commitment
add_content_slide(prs, "Our Commitment: Responsible AI", [
    "• Transparency: Full visibility of our detection algorithms",
    "",
    "• Clinical Oversight: We want NHS guidance, not just approval",
    "",
    "• Data Privacy: GDPR compliant, UK-hosted data",
    "",
    "• Continuous Improvement: Regular review and refinement",
    "",
    "• Not Replacing Humans: AI supports, humans save lives",
    "",
    "We're building technology, but we need clinical wisdom"
])

# Slide 11: The Ask
add_content_slide(prs, "Partnership Request", [
    "Immediate: Clinical review of our safeguarding framework",
    "",
    "Short-term: Guidance on AI response language",
    "",
    "Medium-term: Formal integration with NHS pathways",
    "",
    "Long-term: Research partnership on AI-assisted mental health",
    "",
    "",
    "\"If we can help one veteran get through the night, this is worth it.\""
])

# Slide 12: Thank You
add_title_slide(prs, 
    "Thank You", 
    "Radio Check\nSaving Lives, One Conversation at a Time")

# Save
prs.save('/app/documents/Radio_Check_NHS_Safeguarding_Presentation.pptx')
print("PowerPoint saved to /app/documents/Radio_Check_NHS_Safeguarding_Presentation.pptx")
