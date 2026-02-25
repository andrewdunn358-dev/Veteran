#!/usr/bin/env python3
"""
Generate PDF and PowerPoint versions of the Radio Check Board Presentation
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.util import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import markdown
from weasyprint import HTML, CSS

# Colors - Military/Professional theme
DARK_GREEN = RgbColor(0x1a, 0x3c, 0x34)  # Dark military green
GOLD = RgbColor(0xc9, 0xa2, 0x27)  # Gold accent
WHITE = RgbColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RgbColor(0xf5, 0xf5, 0xf5)

def create_powerpoint():
    """Create PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)
    
    # Helper function to add a title slide
    def add_title_slide(title, subtitle=""):
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Background
        background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        background.fill.solid()
        background.fill.fore_color.rgb = DARK_GREEN
        background.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(28)
            p.font.color.rgb = GOLD
            p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    # Helper function to add content slide
    def add_content_slide(title, bullets, two_column=False):
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Header bar
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
        header.fill.solid()
        header.fill.fore_color.rgb = DARK_GREEN
        header.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # Content
        if two_column and len(bullets) > 4:
            mid = len(bullets) // 2
            left_bullets = bullets[:mid]
            right_bullets = bullets[mid:]
            
            # Left column
            left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5.5))
            tf = left_box.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(left_bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(20)
                p.space_after = Pt(12)
            
            # Right column
            right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(6), Inches(5.5))
            tf = right_box.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(right_bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(20)
                p.space_after = Pt(12)
        else:
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
            tf = content_box.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(22)
                p.space_after = Pt(14)
        
        return slide
    
    # Helper for table slides
    def add_table_slide(title, headers, rows):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Header bar
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
        header.fill.solid()
        header.fill.fore_color.rgb = DARK_GREEN
        header.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # Table
        cols = len(headers)
        table_rows = len(rows) + 1
        table = slide.shapes.add_table(table_rows, cols, Inches(0.5), Inches(1.5), Inches(12.333), Inches(0.5 * table_rows)).table
        
        # Header row
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_GREEN
            p = cell.text_frame.paragraphs[0]
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.size = Pt(16)
        
        # Data rows
        for row_idx, row in enumerate(rows):
            for col_idx, val in enumerate(row):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(val)
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(14)
        
        return slide
    
    # ==================== SLIDES ====================
    
    # Slide 1: Title
    add_title_slide("Radio Check", "Veterans Mental Health & Peer Support Platform\nBoard Presentation - February 2026")
    
    # Slide 2: Executive Summary
    add_content_slide("Executive Summary", [
        "Comprehensive mental health platform for UK military veterans",
        "24/7 access to AI-powered support companions",
        "Peer support networks and professional counselling",
        "4 interconnected platforms working seamlessly",
        "7 unique AI Support Characters",
        "Enterprise-grade security with GDPR/BACP compliance"
    ])
    
    # Slide 3: Platform Components
    add_table_slide("Platform Components", 
        ["Component", "Purpose", "Users"],
        [
            ["Mobile App", "Primary veteran-facing application", "Veterans & families"],
            ["Staff Portal", "Peer supporter & counsellor workspace", "Support staff"],
            ["Admin Portal", "System administration & analytics", "Administrators"],
            ["Marketing Website", "Public information & downloads", "General public"]
        ]
    )
    
    # Slide 4: Live URLs
    add_table_slide("Live Platform URLs",
        ["Platform", "URL", "Hosting"],
        [
            ["Mobile App", "app.radiocheck.me", "Vercel"],
            ["Staff Portal", "staff.radiocheck.me", "20i"],
            ["Admin Portal", "admin.radiocheck.me", "20i"],
            ["Backend API", "veterans-support-api.onrender.com", "Render"]
        ]
    )
    
    # Slide 5: AI Support Characters
    add_content_slide("AI Support Companions (24/7)", [
        "Tommy - Battle Buddy: Straightforward military peer support",
        "Doris - Emotional Support: Nurturing, compassionate listener",
        "Bob - Para Veteran: Real-talk from ex-Parachute Regiment",
        "Finch - Military Law: UK military legal expertise",
        "Margie - Addiction Support: Alcohol, drugs, gambling specialist",
        "Hugo - Wellbeing Coach: Mental health and daily habits",
        "Rita - Family Support: For partners, spouses, loved ones"
    ])
    
    # Slide 6: Mobile App Features
    add_content_slide("Mobile Application Features", [
        "Crisis Support - One-tap access to emergency helplines",
        "AI Chat - 7 specialized support characters",
        "Peer Support - Connect with trained veteran supporters",
        "Callback Requests - Schedule calls from staff",
        "Self-Care Tools - Mood journal, breathing exercises, grounding",
        "Mental Health Screening - PHQ-9 and GAD-7 assessments",
        "Support Organizations Directory",
        "Dark/Light Theme Support"
    ], two_column=True)
    
    # Slide 7: Staff Portal Features
    add_content_slide("Staff Portal Features", [
        "Live Dashboard - Online staff count, metrics",
        "Shift Calendar - Monthly view with staff names",
        "Shift Swap System - Request and approve swaps",
        "Callback Queue - Manage incoming requests",
        "Live Chat - Real-time text support",
        "WebRTC Audio Calls - Peer-to-peer voice calls",
        "Safeguarding Alerts - Risk monitoring dashboard"
    ])
    
    # Slide 8: Admin Portal Features
    add_content_slide("Admin Portal Features", [
        "User Management - Full user directory and roles",
        "Staff Management - Counsellors and peer supporters",
        "Content Management System (CMS) - Edit all app content",
        "Beta Feedback System - Survey management and analytics",
        "Compliance Documents - ROPA, GDPR, BACP downloads",
        "Logs & Analytics - Call logs, chat history, alerts",
        "Organization Directory Management"
    ])
    
    # Slide 9: Technology Stack
    add_table_slide("Technology Stack",
        ["Layer", "Technology", "Rationale"],
        [
            ["Frontend", "React Native (Expo)", "Cross-platform: iOS, Android, Web"],
            ["Backend", "Python FastAPI", "High-performance async API"],
            ["Database", "MongoDB", "Flexible document storage"],
            ["AI Engine", "OpenAI GPT-4", "Intelligent chat companions"],
            ["Email", "Resend", "Transactional email service"],
            ["Auth", "JWT Tokens", "Secure session management"]
        ]
    )
    
    # Slide 10: Security & Compliance
    add_content_slide("Security & Compliance", [
        "GDPR Compliance - Consent management, data export, right to deletion",
        "BACP Clinical Standards - Ethical framework, safeguarding protocols",
        "HTTPS Everywhere - TLS encryption on all endpoints",
        "JWT Authentication - Secure token-based sessions",
        "Password Hashing - bcrypt with salt",
        "Session Timeout - 2-hour inactivity logout",
        "Rate Limiting - IP-based request throttling",
        "Field Encryption - Sensitive data encrypted at rest"
    ], two_column=True)
    
    # Slide 11: Development Investment
    add_table_slide("Development Investment Summary",
        ["Component", "Hours", "Percentage"],
        [
            ["Mobile Application", "548", "33%"],
            ["Admin Portal", "492", "30%"],
            ["Staff Portal", "324", "20%"],
            ["Backend API", "188", "11%"],
            ["Security & Compliance", "48", "3%"],
            ["Marketing Website", "40", "2%"],
            ["TOTAL", "1,640", "100%"]
        ]
    )
    
    # Slide 12: Commercial Value
    add_content_slide("Equivalent Commercial Value", [
        "Total Development Hours: 1,640 hours",
        "Junior Developer Rate (£50/hr): £82,000",
        "Mid-level Developer Rate (£80/hr): £131,200",
        "Senior Developer Rate (£120/hr): £196,800",
        "Blended Rate Estimate (£85/hr): £139,400",
        "",
        "Note: Excludes UI/UX design, project management,",
        "QA testing, hosting, and ongoing maintenance"
    ])
    
    # Slide 13: Future Roadmap
    add_content_slide("Future Roadmap", [
        "Phase 2 (Planned):",
        "  - Push Notifications for appointments",
        "  - Video Call Support (WebRTC)",
        "  - Enhanced Mood Tracking with graphs",
        "  - Welsh Language Support",
        "  - Structured CBT Courses",
        "",
        "Phase 3 (Future):",
        "  - App Store Release (iOS/Android)",
        "  - Offline Mode capabilities",
        "  - NHS Integration",
        "  - Multi-organization support"
    ])
    
    # Slide 14: Thank You
    add_title_slide("Thank You", "Questions?\n\napp.radiocheck.me | admin.radiocheck.me | staff.radiocheck.me")
    
    # Save
    output_path = "/app/admin-site/docs/Radio_Check_Board_Presentation.pptx"
    prs.save(output_path)
    print(f"PowerPoint saved: {output_path}")
    return output_path


def create_pdf():
    """Create PDF from markdown"""
    
    # Read markdown
    with open("/app/board_presentation.md", "r") as f:
        md_content = f.read()
    
    # Convert to HTML
    html_content = markdown.markdown(md_content, extensions=['tables', 'fenced_code'])
    
    # Add styling
    styled_html = f"""
    <!DOCTYPE html>
    <html>
    <head>
        <meta charset="UTF-8">
        <style>
            @page {{
                size: A4;
                margin: 2cm;
            }}
            body {{
                font-family: 'Segoe UI', Arial, sans-serif;
                font-size: 11pt;
                line-height: 1.6;
                color: #333;
            }}
            h1 {{
                color: #1a3c34;
                font-size: 28pt;
                text-align: center;
                border-bottom: 3px solid #c9a227;
                padding-bottom: 10px;
                margin-top: 0;
            }}
            h2 {{
                color: #1a3c34;
                font-size: 18pt;
                border-bottom: 2px solid #1a3c34;
                padding-bottom: 5px;
                margin-top: 30px;
            }}
            h3 {{
                color: #1a3c34;
                font-size: 14pt;
                margin-top: 20px;
            }}
            h4 {{
                color: #555;
                font-size: 12pt;
            }}
            table {{
                width: 100%;
                border-collapse: collapse;
                margin: 15px 0;
                font-size: 10pt;
            }}
            th {{
                background-color: #1a3c34;
                color: white;
                padding: 10px;
                text-align: left;
                font-weight: bold;
            }}
            td {{
                padding: 8px 10px;
                border-bottom: 1px solid #ddd;
            }}
            tr:nth-child(even) {{
                background-color: #f9f9f9;
            }}
            code {{
                background-color: #f4f4f4;
                padding: 2px 6px;
                border-radius: 3px;
                font-family: 'Consolas', monospace;
                font-size: 10pt;
            }}
            pre {{
                background-color: #f4f4f4;
                padding: 15px;
                border-radius: 5px;
                overflow-x: auto;
                font-size: 9pt;
            }}
            strong {{
                color: #1a3c34;
            }}
            hr {{
                border: none;
                border-top: 1px solid #ddd;
                margin: 20px 0;
            }}
            ul, ol {{
                margin-left: 20px;
            }}
            li {{
                margin-bottom: 5px;
            }}
            .cover-info {{
                text-align: center;
                margin: 20px 0;
                color: #666;
            }}
        </style>
    </head>
    <body>
        {html_content}
    </body>
    </html>
    """
    
    # Generate PDF
    output_path = "/app/admin-site/docs/Radio_Check_Board_Presentation.pdf"
    HTML(string=styled_html).write_pdf(output_path)
    print(f"PDF saved: {output_path}")
    return output_path


if __name__ == "__main__":
    print("Creating Board Presentation files...")
    print("-" * 50)
    
    try:
        pdf_path = create_pdf()
        print(f"✓ PDF created successfully")
    except Exception as e:
        print(f"✗ PDF creation failed: {e}")
    
    try:
        pptx_path = create_powerpoint()
        print(f"✓ PowerPoint created successfully")
    except Exception as e:
        print(f"✗ PowerPoint creation failed: {e}")
    
    print("-" * 50)
    print("Done!")
