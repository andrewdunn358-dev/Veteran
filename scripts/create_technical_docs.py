#!/usr/bin/env python3
"""
Generate Technical Overview documents for Radio Check
- PowerPoint (PPTX)
- Word Document (DOCX)
Using brand dark theme colors
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxRGB
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

# Brand Colors
DARK_BG = RGBColor(0x1a, 0x23, 0x32)
SURFACE = RGBColor(0x2d, 0x37, 0x48)
PRIMARY = RGBColor(0x4a, 0x90, 0xe2)
WHITE = RGBColor(0xff, 0xff, 0xff)
TEXT_SECONDARY = RGBColor(0xb0, 0xc4, 0xde)

LOGO_PATH = "/app/admin-site/logo.png"

def create_powerpoint():
    """Create branded Technical Overview PowerPoint"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    def add_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = DARK_BG
        bg.line.fill.background()
    
    def add_title_slide(title, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        
        if os.path.exists(LOGO_PATH):
            slide.shapes.add_picture(LOGO_PATH, Inches(5.9), Inches(1.5), height=Inches(1.5))
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.7), Inches(12.333), Inches(1))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = TEXT_SECONDARY
            p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(title, bullets, small_font=False):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
        header.fill.solid()
        header.fill.fore_color.rgb = SURFACE
        header.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        font_size = Pt(18) if small_font else Pt(22)
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = font_size
            p.font.color.rgb = WHITE
            p.space_after = Pt(10 if small_font else 14)
        
        return slide
    
    def add_table_slide(title, headers, rows, small=False):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
        header.fill.solid()
        header.fill.fore_color.rgb = SURFACE
        header.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        cols = len(headers)
        tbl_rows = len(rows) + 1
        row_height = 0.4 if small else 0.5
        x, y, w, h = Inches(0.7), Inches(1.6), Inches(11.9), Inches(row_height * (len(rows) + 1))
        table = slide.shapes.add_table(tbl_rows, cols, x, y, w, h).table
        
        font_size = Pt(13) if small else Pt(16)
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRIMARY
            p = cell.text_frame.paragraphs[0]
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.size = font_size
        
        data_font = Pt(11) if small else Pt(14)
        for row_idx, row in enumerate(rows):
            for col_idx, val in enumerate(row):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(val)
                cell.fill.solid()
                cell.fill.fore_color.rgb = SURFACE
                p = cell.text_frame.paragraphs[0]
                p.font.color.rgb = WHITE
                p.font.size = data_font
        
        return slide
    
    # ==================== SLIDES ====================
    
    # Slide 1: Title
    add_title_slide("Radio Check", "Technical Overview")
    
    # Slide 2: Architecture Overview
    add_content_slide("System Architecture", [
        "3-tier architecture: Client → API → Database",
        "React Native (Expo) frontend - web platform",
        "FastAPI backend with async Python",
        "MongoDB Atlas for data persistence",
        "OpenAI GPT-4 for AI chat",
        "Resend for transactional email",
        "Socket.IO for real-time features"
    ])
    
    # Slide 3: Technology Stack
    add_table_slide("Technology Stack",
        ["Layer", "Technology", "Purpose"],
        [
            ["Frontend", "React Native (Expo)", "Cross-platform web app"],
            ["Backend", "Python FastAPI", "High-performance async API"],
            ["Database", "MongoDB Atlas", "Document database"],
            ["AI Engine", "OpenAI GPT-4", "Chat companions"],
            ["Email", "Resend", "Notifications"],
            ["Real-time", "Socket.IO", "WebSocket connections"],
            ["Auth", "JWT Tokens", "Secure authentication"]
        ]
    )
    
    # Slide 4: Backend Structure
    add_content_slide("Backend Architecture", [
        "server.py - Main FastAPI application (3000+ lines)",
        "20+ API router modules for separation of concerns",
        "Pydantic models for request/response validation",
        "safety.py - AI safeguarding and risk detection",
        "encryption.py - Field-level AES encryption",
        "cron_runner.py - Scheduled task execution"
    ])
    
    # Slide 5: API Router Modules
    add_content_slide("API Router Modules", [
        "auth.py - Authentication & user management",
        "cms.py - Content management system",
        "shifts.py / shift_swaps.py - Staff rota",
        "callbacks.py - Callback request handling",
        "live_chat.py - Real-time chat (Socket.IO)",
        "safeguarding.py - Risk alerts & screening",
        "surveys.py - Beta feedback system",
        "buddy_finder.py - Peer matching algorithm",
        "resources.py / podcasts.py - Content libraries"
    ], small_font=True)
    
    # Slide 6: API Endpoints Summary
    add_table_slide("Key API Endpoints",
        ["Category", "Example Endpoints", "Methods"],
        [
            ["Authentication", "/auth/login, /auth/register", "POST"],
            ["Staff", "/counsellors, /peer-supporters", "CRUD"],
            ["Shifts", "/shifts, /shift-swaps", "CRUD"],
            ["Callbacks", "/callbacks, /callbacks/{id}/take", "POST, PATCH"],
            ["AI Chat", "/ai/chat, /ai/characters", "POST, GET"],
            ["CMS", "/cms/pages, /cms/sections, /cms/cards", "CRUD"],
            ["Surveys", "/surveys/response, /surveys/stats", "POST, GET"],
            ["Safeguarding", "/safeguarding-alerts, /panic-alert", "GET, POST"]
        ], small=True
    )
    
    # Slide 7: Database Collections
    add_table_slide("Database Schema",
        ["Collection", "Purpose", "Key Fields"],
        [
            ["users", "User accounts", "email, password_hash, role"],
            ["counsellors", "Counsellor profiles", "name, specialization, status"],
            ["peer_supporters", "Peer profiles", "firstName, area, background"],
            ["shifts", "Staff rota", "staff_id, date, start/end_time"],
            ["callbacks", "Callback requests", "name, phone, status, assigned_to"],
            ["safeguarding_alerts", "Risk alerts", "risk_level, session_id, geo_*"],
            ["cms_pages/sections/cards", "CMS content", "slug, title, content"],
            ["survey_responses", "Feedback data", "type, wellbeing, nps"]
        ], small=True
    )
    
    # Slide 8: Security Implementation
    add_content_slide("Security Implementation", [
        "JWT authentication with 24-hour expiry",
        "bcrypt password hashing with salt",
        "Field-level AES-256 encryption (phone, email, name)",
        "Rate limiting: 20 req/min, 5 burst, auto-block",
        "Session timeout after 2 hours inactivity",
        "CORS whitelist for allowed origins",
        "Input validation via Pydantic models"
    ])
    
    # Slide 9: AI Integration
    add_content_slide("AI Integration", [
        "7 AI personas with specialised system prompts",
        "OpenAI GPT-4 model integration",
        "Safety monitoring for suicide/self-harm keywords",
        "Risk scoring: GREEN → YELLOW → AMBER → RED",
        "Automatic crisis resource injection",
        "Geolocation capture on high-risk alerts",
        "Session-level rate limiting (50 messages)"
    ])
    
    # Slide 10: Deployment Infrastructure
    add_table_slide("Deployment Infrastructure",
        ["Service", "Platform", "Purpose"],
        [
            ["Mobile App", "Vercel", "Frontend hosting, CDN"],
            ["Backend API", "Render", "Python app hosting"],
            ["Admin Portal", "20i", "Static site hosting"],
            ["Staff Portal", "20i", "Static site hosting"],
            ["Database", "MongoDB Atlas", "Managed database"]
        ]
    )
    
    # Slide 11: Environment & Config
    add_content_slide("Configuration & Environment", [
        "Backend .env: MONGO_URL, JWT_SECRET_KEY, OPENAI_API_KEY",
        "Frontend .env: REACT_APP_BACKEND_URL",
        "Cron jobs on Render for shift reminders (15 min)",
        "Data retention cleanup daily at 3 AM",
        "Vercel deploy hooks for CI/CD",
        "Static portals require manual FTP upload to 20i"
    ])
    
    # Slide 12: End
    add_title_slide("Radio Check", "Technical Overview v2.6.0")
    
    # Save
    output = "/app/admin-site/docs/Radio_Check_Technical_Overview.pptx"
    prs.save(output)
    print(f"PowerPoint saved: {output}")
    return output


def create_docx():
    """Create branded Technical Overview Word document"""
    doc = Document()
    
    for section in doc.sections:
        section.top_margin = DocxInches(0.75)
        section.bottom_margin = DocxInches(0.75)
        section.left_margin = DocxInches(0.75)
        section.right_margin = DocxInches(0.75)
    
    def add_heading(text, level=1):
        h = doc.add_heading(text, level)
        for run in h.runs:
            run.font.color.rgb = DocxRGB(0x1a, 0x23, 0x32)
        return h
    
    def add_table(headers, rows):
        table = doc.add_table(rows=1, cols=len(headers))
        table.style = 'Table Grid'
        
        hdr_cells = table.rows[0].cells
        for i, h in enumerate(headers):
            hdr_cells[i].text = h
            for p in hdr_cells[i].paragraphs:
                for run in p.runs:
                    run.font.bold = True
                    run.font.color.rgb = DocxRGB(0xff, 0xff, 0xff)
            shading = OxmlElement('w:shd')
            shading.set(qn('w:fill'), '4a90e2')
            hdr_cells[i]._tc.get_or_add_tcPr().append(shading)
        
        for row in rows:
            row_cells = table.add_row().cells
            for i, val in enumerate(row):
                row_cells[i].text = str(val)
        
        doc.add_paragraph()
        return table
    
    def add_code_block(code):
        p = doc.add_paragraph()
        p.style = 'No Spacing'
        run = p.add_run(code)
        run.font.name = 'Consolas'
        run.font.size = DocxPt(9)
        doc.add_paragraph()
    
    # Logo
    if os.path.exists(LOGO_PATH):
        doc.add_picture(LOGO_PATH, width=DocxInches(1.5))
        doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # Title
    title = doc.add_heading('Radio Check', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    subtitle = doc.add_paragraph('Technical Overview')
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle.runs[0].font.size = DocxPt(16)
    subtitle.runs[0].font.color.rgb = DocxRGB(0x4a, 0x90, 0xe2)
    
    doc.add_paragraph()
    
    # 1. Architecture Overview
    add_heading('1. Architecture Overview')
    
    add_heading('System Architecture', level=2)
    doc.add_paragraph(
        'Radio Check uses a 3-tier architecture with a React Native frontend, '
        'FastAPI backend, and MongoDB database. The system integrates with OpenAI '
        'for AI chat and Resend for email notifications.'
    )
    
    arch_bullets = [
        'Client Layer: React Native (Expo) web app, Static HTML portals',
        'API Layer: FastAPI with authentication, rate limiting, CORS',
        'Data Layer: MongoDB Atlas, OpenAI GPT-4, Resend Email'
    ]
    for b in arch_bullets:
        doc.add_paragraph(b, style='List Bullet')
    doc.add_paragraph()
    
    add_heading('Technology Stack', level=2)
    add_table(
        ['Layer', 'Technology', 'Purpose'],
        [
            ['Frontend', 'React Native (Expo) SDK 52', 'Cross-platform web application'],
            ['Backend', 'Python FastAPI 0.100+', 'High-performance async API'],
            ['Database', 'MongoDB 6.0+', 'Document database'],
            ['AI Engine', 'OpenAI GPT-4', 'Chat companions'],
            ['Email', 'Resend', 'Transactional email'],
            ['Real-time', 'Socket.IO 5.x', 'WebSocket connections'],
            ['Authentication', 'JWT', 'Token-based auth']
        ]
    )
    
    # 2. Backend Architecture
    add_heading('2. Backend Architecture')
    
    add_heading('Project Structure', level=2)
    structure = """
/backend
├── server.py           # Main FastAPI application
├── safety.py           # AI safeguarding module
├── encryption.py       # Field-level encryption
├── cron_runner.py      # Scheduled task runner
├── /routers            # API route modules (20+)
└── /models             # Pydantic data models
"""
    add_code_block(structure)
    
    add_heading('Core Modules', level=2)
    modules = [
        'server.py - Main application, CORS, rate limiting, JWT auth, AI personas',
        'safety.py - Risk detection, severity scoring, crisis resource injection',
        'encryption.py - AES-256 field encryption for sensitive data'
    ]
    for m in modules:
        doc.add_paragraph(m, style='List Bullet')
    doc.add_paragraph()
    
    add_heading('API Router Modules', level=2)
    routers = [
        'auth.py - Authentication endpoints',
        'cms.py - Content management',
        'shifts.py - Staff rota management',
        'shift_swaps.py - Shift swap requests',
        'callbacks.py - Callback handling',
        'live_chat.py - Real-time chat (Socket.IO)',
        'safeguarding.py - Risk alerts',
        'surveys.py - Beta feedback',
        'buddy_finder.py - Peer matching',
        'resources.py - Resource library',
        'podcasts.py - Podcast management'
    ]
    for r in routers:
        doc.add_paragraph(r, style='List Bullet')
    doc.add_paragraph()
    
    # 3. API Reference
    add_heading('3. API Reference')
    
    add_heading('Authentication Endpoints', level=2)
    add_table(
        ['Method', 'Endpoint', 'Description'],
        [
            ['POST', '/api/auth/login', 'User login, returns JWT'],
            ['POST', '/api/auth/register', 'Create new user'],
            ['GET', '/api/auth/me', 'Get current user'],
            ['GET', '/api/auth/users', 'List all users (admin)'],
            ['POST', '/api/auth/change-password', 'Change password'],
            ['POST', '/api/auth/reset-password', 'Reset with token']
        ]
    )
    
    add_heading('Staff Management', level=2)
    add_table(
        ['Method', 'Endpoint', 'Description'],
        [
            ['GET/POST', '/api/counsellors', 'List/Create counsellors'],
            ['GET/PUT/DELETE', '/api/counsellors/{id}', 'CRUD operations'],
            ['PATCH', '/api/counsellors/{id}/status', 'Update availability'],
            ['GET/POST', '/api/peer-supporters', 'List/Create peers'],
            ['PATCH', '/api/peer-supporters/{id}/status', 'Update availability']
        ]
    )
    
    add_heading('Shift Management', level=2)
    add_table(
        ['Method', 'Endpoint', 'Description'],
        [
            ['GET/POST', '/api/shifts', 'List/Create shifts'],
            ['GET', '/api/shifts/my-shifts', 'User\'s shifts'],
            ['GET', '/api/shifts/today', 'Today\'s shifts'],
            ['POST', '/api/shift-swaps', 'Request swap'],
            ['PATCH', '/api/shift-swaps/{id}/approve', 'Admin approve']
        ]
    )
    
    add_heading('Callback System', level=2)
    add_table(
        ['Method', 'Endpoint', 'Description'],
        [
            ['POST', '/api/callbacks', 'Create callback request'],
            ['GET', '/api/callbacks', 'List all callbacks'],
            ['PATCH', '/api/callbacks/{id}/take', 'Claim callback'],
            ['PATCH', '/api/callbacks/{id}/complete', 'Mark completed']
        ]
    )
    
    add_heading('Content Management (CMS)', level=2)
    add_table(
        ['Method', 'Endpoint', 'Description'],
        [
            ['GET', '/api/cms/pages', 'List CMS pages'],
            ['GET', '/api/cms/pages/{slug}', 'Get page with sections'],
            ['POST/PUT', '/api/cms/sections', 'Create/Update section'],
            ['POST/PUT', '/api/cms/cards', 'Create/Update card'],
            ['POST', '/api/cms/seed-public', 'Seed default content']
        ]
    )
    
    # 4. Database Schema
    add_heading('4. Database Schema')
    
    add_heading('Collections', level=2)
    add_table(
        ['Collection', 'Purpose', 'Key Fields'],
        [
            ['users', 'User accounts', 'email, password_hash, role, name'],
            ['counsellors', 'Counsellor profiles', 'name, specialization, phone, status'],
            ['peer_supporters', 'Peer profiles', 'firstName, area, background, status'],
            ['shifts', 'Staff rota', 'staff_id, date, start_time, end_time'],
            ['callbacks', 'Callback requests', 'name, phone, message, status'],
            ['safeguarding_alerts', 'Risk alerts', 'risk_level, session_id, geo_*'],
            ['cms_pages', 'CMS pages', 'slug, title, is_visible'],
            ['cms_sections', 'Page sections', 'page_slug, section_type, order'],
            ['cms_cards', 'Content cards', 'section_id, card_type, title'],
            ['survey_responses', 'Feedback', 'type, wellbeing, nps']
        ]
    )
    
    # 5. Security
    add_heading('5. Security Implementation')
    
    add_heading('Authentication', level=2)
    auth_items = [
        'JWT tokens with 24-hour expiry',
        'bcrypt password hashing with salt',
        'Role-based access control (admin, counsellor, peer)'
    ]
    for a in auth_items:
        doc.add_paragraph(a, style='List Bullet')
    doc.add_paragraph()
    
    add_heading('Rate Limiting', level=2)
    rate_items = [
        '20 requests per minute per IP',
        '5 burst requests in 5 seconds triggers warning',
        'Automatic IP blocking for 5 minutes',
        '50 AI messages per session limit'
    ]
    for r in rate_items:
        doc.add_paragraph(r, style='List Bullet')
    doc.add_paragraph()
    
    add_heading('Data Encryption', level=2)
    enc_items = [
        'AES-256-GCM field-level encryption',
        'Encrypted fields: phone, email, name, sms, whatsapp',
        'Automatic encrypt on write, decrypt on read'
    ]
    for e in enc_items:
        doc.add_paragraph(e, style='List Bullet')
    doc.add_paragraph()
    
    # 6. AI Integration
    add_heading('6. AI Integration')
    
    add_heading('Character System', level=2)
    add_table(
        ['Character', 'Focus', 'Prompt Size'],
        [
            ['Tommy', 'General peer support', '~3000 tokens'],
            ['Rachel', 'Criminal justice support', '~3000 tokens'],
            ['Bob', 'Para veteran', '~2500 tokens'],
            ['Finch', 'Legal information', '~2000 tokens'],
            ['Margie', 'Addiction support', '~2500 tokens'],
            ['Hugo', 'Wellbeing coach', '~2000 tokens'],
            ['Rita', 'Family support', '~2000 tokens']
        ]
    )
    
    add_heading('Safety Monitoring', level=2)
    safety_items = [
        'Risk levels: GREEN, YELLOW, AMBER, RED',
        'Keyword detection for suicide/self-harm',
        'Crisis resource injection on high risk',
        'Geolocation capture via IP lookup',
        'Automatic safeguarding alert creation'
    ]
    for s in safety_items:
        doc.add_paragraph(s, style='List Bullet')
    doc.add_paragraph()
    
    # 7. Deployment
    add_heading('7. Deployment Infrastructure')
    
    add_table(
        ['Service', 'Platform', 'Purpose'],
        [
            ['Mobile App', 'Vercel', 'Frontend hosting, CDN'],
            ['Backend API', 'Render', 'Python app hosting'],
            ['Admin Portal', '20i', 'Static site hosting'],
            ['Staff Portal', '20i', 'Static site hosting'],
            ['Database', 'MongoDB Atlas', 'Managed database']
        ]
    )
    
    add_heading('Environment Variables', level=2)
    env_items = [
        'MONGO_URL - MongoDB connection string',
        'DB_NAME - Database name',
        'JWT_SECRET_KEY - Token signing key',
        'OPENAI_API_KEY - AI API key',
        'RESEND_API_KEY - Email API key',
        'ENCRYPTION_KEY - Data encryption key',
        'REACT_APP_BACKEND_URL - Frontend API URL'
    ]
    for e in env_items:
        doc.add_paragraph(e, style='List Bullet')
    doc.add_paragraph()
    
    add_heading('Cron Jobs', level=2)
    cron_items = [
        'Shift reminders - every 15 minutes',
        'Data retention cleanup - daily at 3 AM'
    ]
    for c in cron_items:
        doc.add_paragraph(c, style='List Bullet')
    doc.add_paragraph()
    
    # Footer
    doc.add_paragraph()
    footer = doc.add_paragraph('Radio Check Technical Overview - Version 2.6.0')
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer.runs[0].font.italic = True
    footer.runs[0].font.color.rgb = DocxRGB(0x4a, 0x90, 0xe2)
    
    # Save
    output = "/app/admin-site/docs/Radio_Check_Technical_Overview.docx"
    doc.save(output)
    print(f"Word document saved: {output}")
    return output


if __name__ == "__main__":
    print("Creating Technical Overview documents...")
    print("-" * 50)
    
    try:
        pptx_path = create_powerpoint()
        print(f"✓ PowerPoint created")
    except Exception as e:
        print(f"✗ PowerPoint failed: {e}")
        import traceback
        traceback.print_exc()
    
    try:
        docx_path = create_docx()
        print(f"✓ Word document created")
    except Exception as e:
        print(f"✗ Word document failed: {e}")
        import traceback
        traceback.print_exc()
    
    print("-" * 50)
    print("Done!")
