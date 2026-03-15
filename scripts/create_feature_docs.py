#!/usr/bin/env python3
"""
Generate Feature List documents for Radio Check
- PowerPoint (PPTX)
- Word Document (DOCX)
Using brand dark theme colors
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from docx import Document
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxRGB
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

# Brand Colors (Dark Theme)
DARK_BG = RGBColor(0x1a, 0x23, 0x32)      # #1a2332
SURFACE = RGBColor(0x2d, 0x37, 0x48)       # #2d3748
PRIMARY = RGBColor(0x4a, 0x90, 0xe2)       # #4a90e2
WHITE = RGBColor(0xff, 0xff, 0xff)
TEXT_SECONDARY = RGBColor(0xb0, 0xc4, 0xde) # #b0c4de

LOGO_PATH = "/app/admin-site/logo.png"

def create_powerpoint():
    """Create branded PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    def add_bg(slide):
        """Add dark background"""
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = DARK_BG
        bg.line.fill.background()
    
    def add_title_slide(title, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        
        # Logo
        if os.path.exists(LOGO_PATH):
            slide.shapes.add_picture(LOGO_PATH, Inches(5.9), Inches(1.5), height=Inches(1.5))
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.7), Inches(12.333), Inches(1))
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = TEXT_SECONDARY
            p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(title, bullets):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        
        # Header bar
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
        header.fill.solid()
        header.fill.fore_color.rgb = SURFACE
        header.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(12), Inches(5.3))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {bullet}"
            p.font.size = Pt(22)
            p.font.color.rgb = WHITE
            p.space_after = Pt(14)
        
        return slide
    
    def add_table_slide(title, headers, rows):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide)
        
        # Header bar
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.3))
        header.fill.solid()
        header.fill.fore_color.rgb = SURFACE
        header.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # Table
        cols = len(headers)
        tbl_rows = len(rows) + 1
        x, y, w, h = Inches(0.7), Inches(1.7), Inches(11.9), Inches(0.5 + 0.5 * len(rows))
        table = slide.shapes.add_table(tbl_rows, cols, x, y, w, h).table
        
        # Header row
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRIMARY
            p = cell.text_frame.paragraphs[0]
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.size = Pt(16)
        
        # Data rows
        for row_idx, row in enumerate(rows):
            for col_idx, val in enumerate(row):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(val)
                cell.fill.solid()
                cell.fill.fore_color.rgb = SURFACE
                p = cell.text_frame.paragraphs[0]
                p.font.color.rgb = WHITE
                p.font.size = Pt(14)
        
        return slide
    
    # ==================== SLIDES ====================
    
    # Slide 1: Title
    add_title_slide("Radio Check", "Feature List")
    
    # Slide 2: Platform Overview
    add_content_slide("Platform Overview", [
        "Comprehensive mental health platform for UK military veterans",
        "24/7 AI-powered support companions",
        "Peer support networks",
        "Professional counselling services",
        "Web-based application accessible anywhere"
    ])
    
    # Slide 3: Crisis Support
    add_content_slide("Crisis Support", [
        "One-tap access to emergency helplines",
        "Samaritans integration (116 123)",
        "Veterans UK Helpline (0808 1914 218)",
        "NHS 111 non-emergency support",
        "999 emergency services access",
        "Automatic AI safeguarding detection"
    ])
    
    # Slide 4: AI Support Companions
    add_table_slide("AI Support Companions",
        ["Character", "Focus Area"],
        [
            ["Tommy", "General peer support - squaddie tone"],
            ["Rachel", "Criminal justice - support for veterans in the justice system"],
            ["Bob", "Para veteran - ex-Parachute Regiment"],
            ["Finch", "Military law - UK legal specialist"],
            ["Margie", "Addiction - alcohol, drugs, gambling"],
            ["Hugo", "Wellbeing coach - mental health habits"],
            ["Rita", "Family support - partners & loved ones"]
        ]
    )
    
    # Slide 5: Peer Support Features
    add_content_slide("Peer Support Features", [
        "Talk to a Veteran - connect with trained supporters",
        "Request a Callback - schedule calls from staff",
        "Buddy Finder - match with compatible peers",
        "Live Chat - real-time text support",
        "Voice Calls - peer-to-peer audio support"
    ])
    
    # Slide 6: Self-Care Tools
    add_content_slide("Self-Care Tools", [
        "Mood Journal - daily emotional tracking",
        "Mood Graph - visualise trends over time",
        "Grounding Exercises - 5-4-3-2-1 technique",
        "Breathing Exercises - guided relaxation",
        "Wellness Resources - curated self-help content"
    ])
    
    # Slide 7: Mental Health Screening
    add_content_slide("Mental Health Screening", [
        "PHQ-9 Depression Assessment (9 questions)",
        "GAD-7 Anxiety Assessment (7 questions)",
        "Severity level interpretation",
        "Share results with counsellor option",
        "Automatic staff alerts for high-risk scores"
    ])
    
    # Slide 8: Specialised Support
    add_content_slide("Specialised Support Areas", [
        "Historical Investigations - lawfare support",
        "Criminal Justice - veterans in/leaving prison",
        "Addictions Directory - substance and gambling resources",
        "Family & Friends - resources for loved ones",
        "Regimental Associations - service-specific organisations",
        "Support Organisations Directory"
    ])
    
    # Slide 9: Staff Portal
    add_content_slide("Staff Portal Features", [
        "Live dashboard with status overview",
        "Shift calendar and availability management",
        "Shift swap request system",
        "Callback queue management",
        "Live chat and voice call support",
        "Safeguarding alert dashboard"
    ])
    
    # Slide 10: Admin Portal
    add_content_slide("Admin Portal Features", [
        "User and staff management",
        "Content Management System (CMS)",
        "Beta feedback survey system",
        "Compliance documentation downloads",
        "Logs and analytics dashboard",
        "System health monitoring"
    ])
    
    # Slide 11: Security & Compliance
    add_content_slide("Security & Compliance", [
        "GDPR compliant - consent, data export, deletion",
        "BACP clinical standards - ethical framework",
        "HTTPS encryption on all endpoints",
        "JWT authentication with session timeout",
        "Password hashing (bcrypt)",
        "Rate limiting and input validation",
        "Field-level encryption for sensitive data"
    ])
    
    # Slide 12: Platform URLs
    add_table_slide("Live Platform",
        ["Platform", "URL"],
        [
            ["Mobile App", "app.radiocheck.me"],
            ["Staff Portal", "staff.radiocheck.me"],
            ["Admin Portal", "admin.radiocheck.me"]
        ]
    )
    
    # Slide 13: End
    add_title_slide("Radio Check", "Supporting Veterans 24/7")
    
    # Save
    output = "/app/admin-site/docs/Radio_Check_Features.pptx"
    prs.save(output)
    print(f"PowerPoint saved: {output}")
    return output


def create_docx():
    """Create branded Word document"""
    doc = Document()
    
    # Set page margins
    for section in doc.sections:
        section.top_margin = DocxInches(0.75)
        section.bottom_margin = DocxInches(0.75)
        section.left_margin = DocxInches(0.75)
        section.right_margin = DocxInches(0.75)
    
    def add_heading(text, level=1):
        h = doc.add_heading(text, level)
        for run in h.runs:
            run.font.color.rgb = DocxRGB(0x1a, 0x23, 0x32)
        return h
    
    def add_table(headers, rows):
        table = doc.add_table(rows=1, cols=len(headers))
        table.style = 'Table Grid'
        
        # Header row
        hdr_cells = table.rows[0].cells
        for i, h in enumerate(headers):
            hdr_cells[i].text = h
            for p in hdr_cells[i].paragraphs:
                for run in p.runs:
                    run.font.bold = True
                    run.font.color.rgb = DocxRGB(0xff, 0xff, 0xff)
            # Set background color
            shading = OxmlElement('w:shd')
            shading.set(qn('w:fill'), '4a90e2')
            hdr_cells[i]._tc.get_or_add_tcPr().append(shading)
        
        # Data rows
        for row in rows:
            row_cells = table.add_row().cells
            for i, val in enumerate(row):
                row_cells[i].text = str(val)
        
        doc.add_paragraph()
        return table
    
    # Logo
    if os.path.exists(LOGO_PATH):
        doc.add_picture(LOGO_PATH, width=DocxInches(1.5))
        doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # Title
    title = doc.add_heading('Radio Check', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    subtitle = doc.add_paragraph('Feature List')
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle.runs[0].font.size = DocxPt(16)
    subtitle.runs[0].font.color.rgb = DocxRGB(0x4a, 0x90, 0xe2)
    
    doc.add_paragraph()
    
    # Platform Overview
    add_heading('Platform Overview')
    doc.add_paragraph(
        'Radio Check is a comprehensive mental health and peer support platform designed for UK military veterans, '
        'providing 24/7 access to AI-powered support, peer networks, and professional counselling.'
    )
    doc.add_paragraph()
    
    # Mobile Application
    add_heading('Mobile Application')
    
    add_heading('Crisis Support', level=2)
    bullets = [
        'One-tap access to emergency helplines',
        'Samaritans integration (116 123)',
        'Veterans UK Helpline (0808 1914 218)',
        'NHS 111 non-emergency support',
        '999 emergency services access'
    ]
    for b in bullets:
        doc.add_paragraph(b, style='List Bullet')
    doc.add_paragraph()
    
    add_heading('AI Support Companions', level=2)
    doc.add_paragraph('Available 24/7 with specialised support personas:')
    add_table(
        ['Character', 'Focus Area'],
        [
            ['Tommy', 'General peer support - straightforward squaddie tone'],
            ['Rachel', 'Criminal justice - support for veterans in the justice system'],
            ['Bob', 'Para veteran peer - ex-Parachute Regiment perspective'],
            ['Finch', 'Military law - UK legal information specialist'],
            ['Margie', 'Addiction support - alcohol, drugs, gambling'],
            ['Hugo', 'Wellbeing coach - mental health and daily habits'],
            ['Rita', 'Family support - for partners, spouses, and loved ones']
        ]
    )
    
    ai_features = [
        'Automatic safeguarding detection',
        'Crisis escalation protocols',
        'Conversation session management',
        'Rate limiting protection'
    ]
    doc.add_paragraph('AI Safety Features:', style='Intense Quote')
    for f in ai_features:
        doc.add_paragraph(f, style='List Bullet')
    doc.add_paragraph()
    
    add_heading('Peer Support', level=2)
    peer_features = [
        'Talk to a Veteran - connect with trained supporters',
        'Request a Callback - schedule calls from staff',
        'Buddy Finder - match with compatible peers',
        'Live Chat - real-time text support',
        'Voice Calls - peer-to-peer audio support'
    ]
    for f in peer_features:
        doc.add_paragraph(f, style='List Bullet')
    doc.add_paragraph()
    
    add_heading('Self-Care Tools', level=2)
    selfcare = [
        'Mood Journal - daily emotional tracking',
        'Mood Graph - visualise trends over time',
        'Grounding Exercises - 5-4-3-2-1 technique',
        'Breathing Exercises - guided relaxation',
        'Wellness Resources - curated self-help content'
    ]
    for s in selfcare:
        doc.add_paragraph(s, style='List Bullet')
    doc.add_paragraph()
    
    add_heading('Mental Health Screening', level=2)
    screening = [
        'PHQ-9 Depression Assessment (9 questions)',
        'GAD-7 Anxiety Assessment (7 questions)',
        'Severity level interpretation',
        'Share results with counsellor option',
        'Automatic staff alerts for high-risk scores'
    ]
    for s in screening:
        doc.add_paragraph(s, style='List Bullet')
    doc.add_paragraph()
    
    add_heading('Specialised Support Areas', level=2)
    specialised = [
        'Historical Investigations - lawfare support',
        'Criminal Justice - veterans in/leaving prison',
        'Addictions Directory - substance and gambling resources',
        'Family & Friends - resources for loved ones',
        'Regimental Associations - service-specific organisations'
    ]
    for s in specialised:
        doc.add_paragraph(s, style='List Bullet')
    doc.add_paragraph()
    
    add_heading('Additional Features', level=2)
    additional = [
        'Support Organisations Directory',
        'Recommended Podcasts',
        'Dark/Light Theme Support',
        'Privacy Settings & Data Controls',
        'AI Consent Management (GDPR compliant)'
    ]
    for a in additional:
        doc.add_paragraph(a, style='List Bullet')
    doc.add_paragraph()
    
    # Staff Portal
    add_heading('Staff Portal')
    
    add_heading('Dashboard', level=2)
    dashboard = ['Live status overview', "Today's shifts display", 'Pending callbacks queue', 'Quick statistics']
    for d in dashboard:
        doc.add_paragraph(d, style='List Bullet')
    doc.add_paragraph()
    
    add_heading('Availability Management', level=2)
    availability = [
        'Status toggle (Available/Busy/Offline)',
        'Shift calendar with monthly view',
        'Add availability slots',
        'Create new shifts'
    ]
    for a in availability:
        doc.add_paragraph(a, style='List Bullet')
    doc.add_paragraph()
    
    add_heading('Shift Management', level=2)
    shifts = ['Shift swap requests', 'Manager approval workflow', 'Email notifications', 'Cover request system']
    for s in shifts:
        doc.add_paragraph(s, style='List Bullet')
    doc.add_paragraph()
    
    add_heading('Callback Management', level=2)
    callbacks = [
        'View callback queue',
        'Callback details (user info, reason, urgency)',
        'Take control / claim callbacks',
        'Status updates',
        'Notes system',
        'Release callback option'
    ]
    for c in callbacks:
        doc.add_paragraph(c, style='List Bullet')
    doc.add_paragraph()
    
    add_heading('Communication Tools', level=2)
    comms = ['Live chat interface', 'Voice calls (WebRTC)', 'Message history']
    for c in comms:
        doc.add_paragraph(c, style='List Bullet')
    doc.add_paragraph()
    
    add_heading('Safeguarding', level=2)
    safeguarding = ['Concern logging', 'Alert dashboard', 'Escalation protocols', 'Geolocation tracking for safety']
    for s in safeguarding:
        doc.add_paragraph(s, style='List Bullet')
    doc.add_paragraph()
    
    # Admin Portal
    add_heading('Admin Portal')
    
    add_heading('Dashboard & Analytics', level=2)
    analytics = ['User statistics', 'Staff overview', 'Callback metrics', 'System health monitoring', 'Activity trend charts']
    for a in analytics:
        doc.add_paragraph(a, style='List Bullet')
    doc.add_paragraph()
    
    add_heading('User Management', level=2)
    users = ['User directory (searchable)', 'Role assignment (Admin/Staff/User)', 'Account enable/disable', 'Password reset', 'User creation']
    for u in users:
        doc.add_paragraph(u, style='List Bullet')
    doc.add_paragraph()
    
    add_heading('Content Management System', level=2)
    cms = ['Page management', 'Section editor', 'Card editor (icons, colours, descriptions)', 'Preview system', 'Default content seeding', 'Home page dynamic content']
    for c in cms:
        doc.add_paragraph(c, style='List Bullet')
    doc.add_paragraph()
    
    add_heading('Beta Feedback System', level=2)
    beta = ['Feature flag toggle', 'Pre/Post usage surveys', 'Survey statistics', 'NPS score tracking', 'Wellbeing metrics', 'CSV data export']
    for b in beta:
        doc.add_paragraph(b, style='List Bullet')
    doc.add_paragraph()
    
    add_heading('Compliance Documentation', level=2)
    compliance_docs = [
        'ROPA (Record of Processing Activities)',
        'GDPR Audit Report',
        'BACP Compliance Framework',
        'Incident Response Plan',
        'Security Review Schedule',
        'Safeguarding Disclaimer',
        'PDF downloads'
    ]
    for c in compliance_docs:
        doc.add_paragraph(c, style='List Bullet')
    doc.add_paragraph()
    
    add_heading('Logs & Monitoring', level=2)
    logs = ['Call intent logs', 'Chat room logs', 'Safeguarding alerts', 'Screening submissions', 'Panic alert logs', 'CSV export for all log types']
    for l in logs:
        doc.add_paragraph(l, style='List Bullet')
    doc.add_paragraph()
    
    # Security & Compliance
    add_heading('Security & Compliance')
    
    add_heading('Data Protection (GDPR)', level=2)
    gdpr = ['Consent management', 'Data export capability', 'Right to deletion', 'Privacy by design', 'Data minimisation']
    for g in gdpr:
        doc.add_paragraph(g, style='List Bullet')
    doc.add_paragraph()
    
    add_heading('Clinical Standards (BACP)', level=2)
    bacp = ['Ethical framework compliance', 'Safeguarding protocols', 'Confidentiality measures', 'Professional boundaries']
    for b in bacp:
        doc.add_paragraph(b, style='List Bullet')
    doc.add_paragraph()
    
    add_heading('Security Measures', level=2)
    security = [
        'HTTPS encryption',
        'JWT authentication',
        'Password hashing (bcrypt)',
        'Session timeout (2 hours)',
        'Rate limiting',
        'Input validation',
        'Field-level encryption'
    ]
    for s in security:
        doc.add_paragraph(s, style='List Bullet')
    doc.add_paragraph()
    
    # Technology Stack
    add_heading('Technology Stack')
    add_table(
        ['Component', 'Technology'],
        [
            ['Frontend', 'React Native (Expo) - Web'],
            ['Backend', 'Python FastAPI'],
            ['Database', 'MongoDB'],
            ['AI Engine', 'OpenAI GPT-4'],
            ['Email', 'Resend'],
            ['Authentication', 'JWT Tokens']
        ]
    )
    
    # Live Platform
    add_heading('Live Platform URLs')
    add_table(
        ['Platform', 'URL'],
        [
            ['Mobile App', 'app.radiocheck.me'],
            ['Staff Portal', 'staff.radiocheck.me'],
            ['Admin Portal', 'admin.radiocheck.me']
        ]
    )
    
    # Footer
    doc.add_paragraph()
    footer = doc.add_paragraph('Radio Check - Supporting Veterans 24/7')
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer.runs[0].font.italic = True
    footer.runs[0].font.color.rgb = DocxRGB(0x4a, 0x90, 0xe2)
    
    # Save
    output = "/app/admin-site/docs/Radio_Check_Features.docx"
    doc.save(output)
    print(f"Word document saved: {output}")
    return output


if __name__ == "__main__":
    print("Creating Feature List documents...")
    print("-" * 50)
    
    try:
        pptx_path = create_powerpoint()
        print(f"✓ PowerPoint created")
    except Exception as e:
        print(f"✗ PowerPoint failed: {e}")
        import traceback
        traceback.print_exc()
    
    try:
        docx_path = create_docx()
        print(f"✓ Word document created")
    except Exception as e:
        print(f"✗ Word document failed: {e}")
        import traceback
        traceback.print_exc()
    
    print("-" * 50)
    print("Done!")
